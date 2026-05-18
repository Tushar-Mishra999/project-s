"""Text extraction from uploaded files — mirrors lib/extract.js."""
import re
import zipfile
import io
from xml.etree import ElementTree as ET


def _detect_headings(lines: list[str]) -> list[dict]:
    headings = []
    for i, line in enumerate(lines):
        stripped = line.strip()
        if not stripped:
            continue
        # Markdown headings
        m = re.match(r'^(#{1,6})\s+(.+)', stripped)
        if m:
            headings.append({"line": i, "level": len(m.group(1)), "text": m.group(2)})
            continue
        # Numbered sections 1. / 1.1 / 2.3.4
        if re.match(r'^\d+(\.\d+)*\.?\s+\S', stripped):
            headings.append({"line": i, "level": 2, "text": stripped})
            continue
        # ALL CAPS lines
        if len(stripped) >= 4 and stripped == stripped.upper() and not stripped.endswith(('.', '?', '!')):
            headings.append({"line": i, "level": 2, "text": stripped})
            continue
        # Title-case followed by blank line
        if (
            i + 1 < len(lines)
            and not lines[i + 1].strip()
            and stripped == stripped.title()
            and len(stripped.split()) <= 8
        ):
            headings.append({"line": i, "level": 3, "text": stripped})
    return headings


def extract_text(file_bytes: bytes, filetype: str) -> dict:
    ft = filetype.lower()

    if "pdf" in ft:
        return _extract_pdf(file_bytes)
    if "docx" in ft or "wordprocessingml" in ft or "openxmlformats-officedocument.wordprocessingml" in ft:
        return _extract_docx(file_bytes)
    if "pptx" in ft or "presentationml" in ft:
        return _extract_pptx(file_bytes)
    if "xlsx" in ft or "spreadsheetml" in ft or "excel" in ft:
        return _extract_xlsx(file_bytes)
    # Plain text fallback
    return _extract_text(file_bytes)


def _extract_pdf(data: bytes) -> dict:
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(stream=data, filetype="pdf")
        pages = [page.get_text() for page in doc]
        text = "\n\n".join(pages)
        lines = text.splitlines()
        return {"text": text, "headings": _detect_headings(lines)}
    except Exception as e:
        return {"text": "", "headings": [], "error": str(e)}


def _extract_docx(data: bytes) -> dict:
    try:
        from docx import Document
        doc = Document(io.BytesIO(data))
        text = "\n".join(p.text for p in doc.paragraphs)
        lines = text.splitlines()
        return {"text": text, "headings": _detect_headings(lines)}
    except Exception as e:
        # Try raw XML extraction as fallback for malformed docx
        try:
            with zipfile.ZipFile(io.BytesIO(data)) as z:
                if "word/document.xml" in z.namelist():
                    xml = z.read("word/document.xml")
                    root = ET.fromstring(xml)
                    texts = [t.text or "" for t in root.iter("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}t")]
                    text = " ".join(texts).strip()
                    return {"text": text, "headings": _detect_headings(text.splitlines())}
        except Exception:
            pass
        return {"text": "", "headings": [], "error": str(e)}


def _extract_pptx(data: bytes) -> dict:
    slides = []
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(data))
        for slide in prs.slides:
            parts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    parts.append(shape.text_frame.text)
            slides.append(" ".join(parts).strip())
    except Exception:
        # Fallback: raw XML extraction
        try:
            with zipfile.ZipFile(io.BytesIO(data)) as z:
                names = sorted(n for n in z.namelist() if re.match(r'ppt/slides/slide\d+\.xml', n))
                for name in names:
                    xml = z.read(name)
                    root = ET.fromstring(xml)
                    ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
                    texts = [t.text or '' for t in root.iter('{http://schemas.openxmlformats.org/drawingml/2006/main}t')]
                    slides.append(" ".join(texts).strip())
        except Exception:
            pass

    text = "\n\n".join(s for s in slides if s)
    return {"text": text, "headings": [], "slides": slides}


def _extract_xlsx(data: bytes) -> dict:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(io.BytesIO(data), data_only=True)
        sheets = []
        for sheet in wb.worksheets:
            rows = []
            for row in sheet.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(cells):
                    rows.append(",".join(cells))
            if rows:
                sheets.append(f"Sheet: {sheet.title}\n" + "\n".join(rows))
        text = "\n\n".join(sheets)
        return {"text": text, "headings": []}
    except Exception as e:
        return {"text": "", "headings": [], "error": str(e)}


def _extract_text(data: bytes) -> dict:
    try:
        text = data.decode("utf-8", errors="replace")
    except Exception:
        text = ""
    lines = text.splitlines()
    return {"text": text, "headings": _detect_headings(lines)}
